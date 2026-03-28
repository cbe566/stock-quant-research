#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
海外對衝基金投資者教育PPT生成腳本
面向中國投資者的對衝基金知識普及
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 全域設定 ====================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 配色方案 — 深藍金融風格
COLOR_DARK_BG = RGBColor(0x0B, 0x1D, 0x3A)       # 深藍背景
COLOR_MID_BG = RGBColor(0x12, 0x2B, 0x4F)         # 中藍背景
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xD6)          # 亮藍強調
COLOR_GOLD = RGBColor(0xD4, 0xA5, 0x37)            # 金色
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)            # 白色
COLOR_LIGHT_GRAY = RGBColor(0xC0, 0xC8, 0xD4)      # 淺灰
COLOR_GREEN = RGBColor(0x2E, 0xCC, 0x71)            # 綠色（正面）
COLOR_RED = RGBColor(0xE7, 0x4C, 0x3C)              # 紅色（警示）
COLOR_ORANGE = RGBColor(0xF3, 0x9C, 0x12)           # 橙色
COLOR_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)        # 深色文字

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ==================== 輔助函數 ====================

def add_bg(slide, color=COLOR_DARK_BG):
    """添加純色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """添加矩形色塊"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Microsoft YaHei"):
    """添加文字框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=16,
                   default_color=COLOR_WHITE, line_spacing=1.3):
    """添加多行文字，每行可自訂格式
    lines: list of dict, 每個 dict 有 text, size(可選), color(可選), bold(可選), bullet(可選)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            line_info = {"text": line_info}

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = line_info.get("text", "")
        if line_info.get("bullet"):
            text = "  •  " + text

        p.text = text
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", False)
        p.font.name = "Microsoft YaHei"
        p.alignment = line_info.get("alignment", PP_ALIGN.LEFT)
        p.space_after = Pt(line_info.get("space_after", 4))
        if line_spacing:
            p.line_spacing = Pt(line_info.get("size", default_size) * line_spacing)

    return txBox

def add_divider_line(slide, left, top, width, color=COLOR_GOLD, height=Pt(2)):
    """添加分隔線"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, body_lines,
             title_color=COLOR_GOLD, bg_color=COLOR_MID_BG, title_size=18, body_size=14):
    """添加卡片元素"""
    card = add_shape_bg(slide, left, top, width, height, bg_color)

    # 標題
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.5),
                 title, font_size=title_size, color=title_color, bold=True)

    # 分隔線
    add_divider_line(slide, left + Inches(0.2), top + Inches(0.6),
                     width - Inches(0.4), color=COLOR_ACCENT, height=Pt(1))

    # 內容
    add_multi_text(slide, left + Inches(0.2), top + Inches(0.7),
                   width - Inches(0.4), height - Inches(0.9),
                   body_lines, default_size=body_size)

    return card

def add_table_slide(slide, left, top, width, height, headers, rows,
                    header_color=COLOR_DARK_BG, header_text_color=COLOR_WHITE,
                    row_color1=RGBColor(0xF5, 0xF7, 0xFA), row_color2=COLOR_WHITE,
                    font_size=12):
    """在深色背景上添加表格"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # 設定表頭
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = header_text_color
            paragraph.font.bold = True
            paragraph.font.name = "Microsoft YaHei"
            paragraph.alignment = PP_ALIGN.CENTER

    # 設定行數據
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color1 if i % 2 == 0 else row_color2
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = COLOR_DARK_TEXT
                paragraph.font.name = "Microsoft YaHei"
                paragraph.alignment = PP_ALIGN.CENTER

    return table_shape

def make_section_title_slide(title, subtitle=""):
    """創建章節標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    add_bg(slide, COLOR_DARK_BG)

    # 裝飾線條
    add_shape_bg(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Pt(3), COLOR_GOLD)

    # 標題
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
                 title, font_size=44, color=COLOR_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(1),
                     subtitle, font_size=22, color=COLOR_LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    return slide


# ==================== 第1頁：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

# 頂部裝飾線
add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), COLOR_GOLD)

# 主標題
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "海外对冲基金", font_size=56, color=COLOR_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "OVERSEAS HEDGE FUNDS", font_size=28, color=COLOR_GOLD, bold=False,
             alignment=PP_ALIGN.CENTER)

# 分隔線
add_divider_line(slide, Inches(4), Inches(4.0), Inches(5.3), COLOR_ACCENT, Pt(2))

# 副標題
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.8),
             "投资者教育系列  |  全面解析全球对冲基金行业", font_size=22,
             color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 底部信息
add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
             "数据来源：HFR · Preqin · Bloomberg · J.P. Morgan · Institutional Investor",
             font_size=14, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10), Inches(0.5),
             "2026年3月", font_size=16, color=COLOR_GOLD, alignment=PP_ALIGN.CENTER)

# 底部裝飾線
add_shape_bg(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, Pt(4), COLOR_GOLD)


# ==================== 第2頁：目錄 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.8),
             "目  录  CONTENTS", font_size=36, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.2), Inches(3), COLOR_ACCENT)

toc_items = [
    ("01", "什么是对冲基金？"),
    ("02", "对冲基金 vs 公募基金 vs 私募股权"),
    ("03", "对冲基金发展史（1949-2026）"),
    ("04", "六大核心策略详解"),
    ("05", "传奇案例与经典交易"),
    ("06", "费用结构：2/20 模式"),
    ("07", "风险全景与全球监管"),
    ("08", "全球市场规模与趋势"),
    ("09", "中国投资者参与渠道"),
    ("10", "核心术语表与总结"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.6) + Inches(0.55) * i
    add_text_box(slide, Inches(1.2), y, Inches(0.8), Inches(0.5),
                 num, font_size=26, color=COLOR_GOLD, bold=True)
    add_text_box(slide, Inches(2.2), y + Inches(0.03), Inches(9), Inches(0.5),
                 title, font_size=20, color=COLOR_WHITE)
    if i < len(toc_items) - 1:
        add_divider_line(slide, Inches(1.2), y + Inches(0.48), Inches(10.5),
                         RGBColor(0x20, 0x40, 0x70), Pt(1))


# ==================== 第3頁：什麼是對衝基金 ====================
make_section_title_slide("01  什么是对冲基金？",
                         "Hedge Fund — 以绝对回报为目标的私募投资基金")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "什么是对冲基金？", font_size=32, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

# 定義
add_multi_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(1.5), [
    {"text": "对冲基金（Hedge Fund）是一种以「绝对回报」为目标的私募投资基金，",
     "size": 16, "color": COLOR_WHITE},
    {"text": "通过多元化策略（做多、做空、杠杆、衍生品等）在各种市场环境下追求正收益。",
     "size": 16, "color": COLOR_WHITE},
])

# 五大核心特徵卡片
features = [
    ("绝对收益导向", "不以跑赢指数为目标\n牛熊市皆追求正回报"),
    ("策略灵活", "可做多、做空\n使用杠杆与衍生品"),
    ("高门槛准入", "仅面向合格投资者\n通常$100万美元以上"),
    ("流动性限制", "设有锁定期\n季度或年度赎回"),
    ("有限合伙制", "LP/GP结构运作\n管理人自有资金投入"),
]

for i, (title, desc) in enumerate(features):
    left = Inches(0.8) + Inches(2.4) * i
    card_shape = add_shape_bg(slide, left, Inches(3.0), Inches(2.2), Inches(2.8), COLOR_MID_BG)
    add_text_box(slide, left + Inches(0.15), Inches(3.1), Inches(1.9), Inches(0.5),
                 title, font_size=16, color=COLOR_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_divider_line(slide, left + Inches(0.3), Inches(3.55), Inches(1.6), COLOR_ACCENT, Pt(1))
    add_text_box(slide, left + Inches(0.15), Inches(3.7), Inches(1.9), Inches(1.8),
                 desc, font_size=13, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Alpha/Beta 公式
add_shape_bg(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.0), COLOR_MID_BG)
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11), Inches(0.8),
             '核心逻辑：  投资组合收益 = Alpha（选股能力）+ Beta（市场风险）× 市场收益      '
             '→  对冲基金目标：最大化 Alpha，降低或消除 Beta',
             font_size=16, color=COLOR_WHITE, bold=False, alignment=PP_ALIGN.CENTER)


# ==================== 第5頁：三方比較 ====================
make_section_title_slide("02  三种基金类型比较",
                         "对冲基金 vs 公募基金 vs 私募股权基金")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))  # 淺色背景方便看表格

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "对冲基金 vs 公募基金 vs 私募股权基金", font_size=28,
             color=COLOR_DARK_BG, bold=True)

headers = ["比较维度", "对冲基金 Hedge Fund", "公募基金 Mutual Fund", "私募股权 Private Equity"]
rows = [
    ["投资标的", "全资产类别（股票、债券、\n商品、外汇、衍生品）", "主要为股票、债券", "非上市公司股权"],
    ["策略", "做多做空、杠杆、套利、\n全球宏观等", "主要做多（Long-only）", "收购、重组、\n改善经营后退出"],
    ["投资者门槛", "高（$100万~$500万）", "低（任何人均可）", "极高（$500万~$1000万）"],
    ["流动性", "较低（季度/年度赎回）", "高（每日可赎回）", "极低（5-10年锁定期）"],
    ["费用结构", "2%管理费+20%业绩报酬", "1%-2%管理费", "2%管理费+20%附带收益"],
    ["收益目标", "绝对收益", "相对收益（跑赢基准）", "高回报（经营改善）"],
    ["监管程度", "较低（2010后加强）", "严格监管", "较低"],
    ["信息透明度", "较低", "高（定期披露持仓）", "较低"],
]

add_table_slide(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8),
                headers, rows, font_size=12,
                header_color=COLOR_DARK_BG)


# ==================== 第7頁：發展史（上） ====================
make_section_title_slide("03  对冲基金发展史",
                         "从1949年到2026年 — 七十余年的传奇历程")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "发展史（上）：创始与崛起  1949-1999", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

# 時間線事件
events_upper = [
    ("1949", "Alfred W. Jones 创立第一支对冲基金\n初始资金10万美元，发明「做多+做空」组合与2/20收费模式", COLOR_ACCENT),
    ("1966", "《财富》杂志首次使用「Hedge Fund」一词\nJones基金5年跑赢最优秀公募基金44%，引爆行业关注", COLOR_GREEN),
    ("1969", "索罗斯创立量子基金 Quantum Fund\nSteinhardt创立Steinhardt Partners", COLOR_ACCENT),
    ("1970s", "第一次行业洗牌 — 石油危机引发股市暴跌\n过度杠杆基金遭受重创，行业规模大幅收缩", COLOR_RED),
    ("1990", "Ken Griffin 22岁创立 Citadel 城堡投资\n在哈佛宿舍安装卫星天线开始交易的传奇故事", COLOR_ACCENT),
    ("1992", "索罗斯做空英镑，获利超10亿英镑\n被称为「击垮英格兰银行的人」", COLOR_GOLD),
    ("1998", "LTCM崩盘 — 诺贝尔奖得主也无法避免的灾难\n美联储协调14家机构注资36.5亿美元救助", COLOR_RED),
    ("1999", "行业增长至~4,000支基金\n管理资产超过4,500亿美元", COLOR_GREEN),
]

for i, (year, desc, color) in enumerate(events_upper):
    y = Inches(1.3) + Inches(0.73) * i
    # 年份
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.5),
                 year, font_size=18, color=color, bold=True)
    # 圓點
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.15), y + Inches(0.08), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    # 描述
    add_text_box(slide, Inches(2.5), y - Inches(0.05), Inches(9.5), Inches(0.7),
                 desc, font_size=13, color=COLOR_LIGHT_GRAY)

# 時間線豎線
add_shape_bg(slide, Inches(2.2), Inches(1.3), Pt(2), Inches(5.6), COLOR_ACCENT)


# ==================== 第9頁：發展史（下） ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "发展史（下）：机构化与新纪元  2000-2026", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

events_lower = [
    ("2000-02", "科网泡沫破裂 — 对冲基金因可做空而表现优异\n养老金、主权财富基金开始大规模配置，行业进入「机构化」时代", COLOR_GREEN),
    ("2005", "行业规模首次突破 1 万亿美元", COLOR_ACCENT),
    ("2007", "巅峰时刻 — 管理资产接近 2 万亿美元\n全球对冲基金数量超过 10,000 支", COLOR_GOLD),
    ("2008", "全球金融危机 — 平均亏损15%-20%\n麦道夫650亿美元骗局曝光，行业信誉遭重创", COLOR_RED),
    ("2010", "《Dodd-Frank法案》通过 — 要求向SEC注册\n引入Form PF报告制度，加强系统性风险监控", COLOR_ACCENT),
    ("2021", "GameStop事件 — 散户对抗对冲基金\nMelvin Capital遭受重创，引发市场公平性讨论", COLOR_ORANGE),
    ("2022", "Citadel创造160亿美元年利润\n单一对冲基金年度最高纪录", COLOR_GREEN),
    ("2025", "历史性里程碑 — 全球管理资产突破 5 万亿美元\n连续两年双位数回报，行业平均回报率11.2%", COLOR_GOLD),
]

for i, (year, desc, color) in enumerate(events_lower):
    y = Inches(1.3) + Inches(0.73) * i
    add_text_box(slide, Inches(0.8), y, Inches(1.2), Inches(0.5),
                 year, font_size=18, color=color, bold=True)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.15), y + Inches(0.08), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(2.5), y - Inches(0.05), Inches(9.5), Inches(0.7),
                 desc, font_size=13, color=COLOR_LIGHT_GRAY)

add_shape_bg(slide, Inches(2.2), Inches(1.3), Pt(2), Inches(5.6), COLOR_ACCENT)


# ==================== 第10頁：六大策略總覽 ====================
make_section_title_slide("04  六大核心策略详解",
                         "股票多空 · 全球宏观 · 事件驱动 · 相对价值 · 量化策略 · 多策略")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "六大核心策略一览", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(3.5), COLOR_ACCENT)

strategies = [
    ("股票多空", "Equity L/S", "同时做多低估股票\n做空高估股票",
     "占比~30%\n2025回报17.6%", "Tiger Global\nLone Pine", COLOR_ACCENT),
    ("全球宏观", "Global Macro", "基于宏观经济趋势\n跨市场方向性投资",
     "投资范围最广\n2025回报7.16%", "桥水基金\n索罗斯", COLOR_GOLD),
    ("事件驱动", "Event-Driven", "利用并购、重组等\n企业事件获利",
     "含并购套利\n困境证券等", "Elliott\nPaulson", COLOR_GREEN),
    ("相对价值", "Relative Value", "利用相关证券间\n的价格偏差套利",
     "低净敞口\n收益稳定", "Citadel\nLTCM(历史)", COLOR_ORANGE),
    ("量化策略", "Quantitative", "完全依赖数学模型\nAI与机器学习驱动",
     "增长最快\nCAGR 11.63%", "文艺复兴\nTwo Sigma", COLOR_ACCENT),
    ("多策略", "Multi-Strategy", "同一基金内配置\n多种策略动态调整",
     "2025年成为\n最大策略类别", "Millennium\nCitadel", RGBColor(0xBB, 0x86, 0xFC)),
]

for i, (name_cn, name_en, logic, data, rep, color) in enumerate(strategies):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + Inches(4.1) * col
    top = Inches(1.3) + Inches(3.0) * row

    # 卡片背景
    add_shape_bg(slide, left, top, Inches(3.8), Inches(2.7), COLOR_MID_BG)

    # 策略名稱
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.5), Inches(0.4),
                 f"{name_cn}  {name_en}", font_size=16, color=color, bold=True)
    add_divider_line(slide, left + Inches(0.15), top + Inches(0.5), Inches(3.5), color, Pt(1))

    # 策略邏輯
    add_text_box(slide, left + Inches(0.15), top + Inches(0.6), Inches(1.7), Inches(0.8),
                 "策略逻辑：", font_size=11, color=COLOR_GOLD, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.85), Inches(1.7), Inches(0.8),
                 logic, font_size=12, color=COLOR_LIGHT_GRAY)

    # 數據
    add_text_box(slide, left + Inches(1.9), top + Inches(0.6), Inches(1.7), Inches(0.8),
                 "关键数据：", font_size=11, color=COLOR_GOLD, bold=True)
    add_text_box(slide, left + Inches(1.9), top + Inches(0.85), Inches(1.7), Inches(0.8),
                 data, font_size=12, color=COLOR_LIGHT_GRAY)

    # 代表基金
    add_text_box(slide, left + Inches(0.15), top + Inches(1.9), Inches(3.5), Inches(0.7),
                 f"代表基金：{rep.replace(chr(10), '、')}", font_size=12, color=COLOR_WHITE)


# ==================== 策略詳解頁面 ====================

# 股票多空策略詳解
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "策略详解 ①  股票多空策略", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

# 左側：策略說明
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5),
         "股票多空 Equity Long/Short", [
    {"text": "核心逻辑", "size": 16, "color": COLOR_GOLD, "bold": True, "space_after": 2},
    {"text": "同时做多被低估的股票，做空被高估的股票", "size": 14, "bullet": True},
    {"text": "通过「对冲」降低市场系统性风险", "size": 14, "bullet": True},
    {"text": "收益来源是选股能力（Alpha），而非市场涨跌（Beta）", "size": 14, "bullet": True},
    {"text": "", "size": 8},
    {"text": "简单例子", "size": 16, "color": COLOR_GOLD, "bold": True, "space_after": 2},
    {"text": "基金经理看好茅台（做多100万），看空某劣质白酒（做空100万）", "size": 14, "bullet": True},
    {"text": "若白酒行业整体下跌10%：", "size": 14, "bullet": True},
    {"text": "  茅台因品质优秀只跌5% → 亏5万", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "  劣质白酒跌15% → 做空赚15万", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "  净收益 = +10万（即使行业下跌也能赚钱）", "size": 14, "color": COLOR_GREEN, "bold": True},
], title_size=16)

# 右側：數據
add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(2.5),
         "关键数据", [
    {"text": "占比：约占对冲基金策略的 30%，最普遍", "size": 14, "bullet": True},
    {"text": "2025年回报率：17.6%（所有策略中表现最佳）", "size": 14, "bullet": True, "color": COLOR_GREEN},
    {"text": "代表基金：Tiger Global、Lone Pine、Viking Global", "size": 14, "bullet": True},
], title_size=16)

add_card(slide, Inches(7.0), Inches(4.1), Inches(5.5), Inches(2.7),
         "净敞口与总敞口", [
    {"text": "总敞口 = |多头| + |空头|", "size": 14, "bullet": True},
    {"text": "例：做多$100万 + 做空$80万 → 总敞口$180万", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "净敞口 = 多头 - 空头", "size": 14, "bullet": True},
    {"text": "例：$100万 - $80万 = $20万（净多头20%）", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "净敞口越低，对冲程度越高，市场风险越小", "size": 14, "color": COLOR_GOLD, "bullet": True},
], title_size=16)


# 全球宏觀策略詳解
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "策略详解 ②  全球宏观策略", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5),
         "全球宏观 Global Macro", [
    {"text": "核心逻辑", "size": 16, "color": COLOR_GOLD, "bold": True, "space_after": 2},
    {"text": "基于对全球宏观经济趋势的判断进行投资", "size": 14, "bullet": True},
    {"text": "可跨市场、跨资产类别、跨地区操作", "size": 14, "bullet": True},
    {"text": "投资范围最广：股票、债券、外汇、商品、利率", "size": 14, "bullet": True},
    {"text": "", "size": 8},
    {"text": "关注因素", "size": 16, "color": COLOR_GOLD, "bold": True, "space_after": 2},
    {"text": "通胀走势与央行利率政策", "size": 14, "bullet": True},
    {"text": "GDP增长与经济周期判断", "size": 14, "bullet": True},
    {"text": "地缘政治事件（战争、贸易摩擦等）", "size": 14, "bullet": True},
    {"text": "汇率走势与资本流动", "size": 14, "bullet": True},
    {"text": "", "size": 8},
    {"text": "经典案例：索罗斯1992年做空英镑 → 获利超10亿英镑", "size": 14, "color": COLOR_GREEN},
], title_size=16)

add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
         "全球宏观策略的思维方式", [
    {"text": "第一步：判断全球经济处于什么阶段", "size": 14, "color": COLOR_ACCENT, "bold": True},
    {"text": "扩张？高峰？收缩？谷底？", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "第二步：判断各国央行政策方向", "size": 14, "color": COLOR_ACCENT, "bold": True},
    {"text": "美联储加息还是降息？欧央行、日银政策？", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "第三步：找到价格偏离基本面的资产", "size": 14, "color": COLOR_ACCENT, "bold": True},
    {"text": "汇率被高估？利率定价不合理？", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "第四步：建立方向性头寸", "size": 14, "color": COLOR_ACCENT, "bold": True},
    {"text": "做多或做空，通常使用杠杆放大收益", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "代表基金：桥水 Bridgewater、Brevan Howard", "size": 14, "color": COLOR_GOLD},
    {"text": "2025年回报率：7.16%", "size": 14, "color": COLOR_GREEN},
], title_size=16)


# ==================== 傳奇案例 ====================
make_section_title_slide("05  传奇案例与经典交易",
                         "桥水 · 文艺复兴 · 城堡 · 索罗斯 · LTCM")

# 橋水基金
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "案例 ①  桥水基金 Bridgewater Associates", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(5), COLOR_ACCENT)

# 左：簡介
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(3.0),
         "基金概况", [
    {"text": "创始人：Ray Dalio（瑞·达利欧）", "size": 14, "bullet": True},
    {"text": "成立时间：1975年", "size": 14, "bullet": True},
    {"text": "管理资产：约 $1,500亿美元", "size": 14, "bullet": True, "color": COLOR_GREEN},
    {"text": "累计净利润：$497亿美元", "size": 14, "bullet": True, "color": COLOR_GREEN},
    {"text": "核心策略：系统化宏观策略", "size": 14, "bullet": True},
    {"text": "总部：美国康涅狄格州", "size": 14, "bullet": True},
], title_size=16)

add_card(slide, Inches(0.8), Inches(4.6), Inches(5.8), Inches(2.5),
         "核心理念", [
    {"text": "「原则」— 极度透明、系统化决策、从错误中学习", "size": 14, "bullet": True},
    {"text": "经济机器理论 — 将经济视为一台可分析的机器", "size": 14, "bullet": True},
    {"text": "风险平价 — 根据风险（而非金额）平均分配资产", "size": 14, "bullet": True},
], title_size=16)

# 右：旗艦基金
add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.8),
         "两大旗舰基金", [
    {"text": "Pure Alpha 纯阿尔法", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "主动宏观策略，追求绝对收益", "size": 14, "bullet": True},
    {"text": "根据全球经济指标做出方向性判断", "size": 14, "bullet": True},
    {"text": "灵活配置股票、债券、外汇、商品", "size": 14, "bullet": True},
    {"text": "", "size": 8},
    {"text": "All Weather 全天候", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "风险平价（Risk Parity）策略", "size": 14, "bullet": True},
    {"text": "将资产按四种经济环境分配：", "size": 14, "bullet": True},
    {"text": "  增长上升 → 股票、商品", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "  增长下降 → 债券、通胀连结债", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "  通胀上升 → 商品、通胀连结债", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "  通胀下降 → 股票、债券", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "目标：在任何经济环境下都能获得合理回报", "size": 14, "color": COLOR_GOLD},
], title_size=16)


# 文藝復興科技
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "案例 ②  文艺复兴科技 Renaissance Technologies", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(5), COLOR_ACCENT)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5),
         "基金概况", [
    {"text": "创始人：Jim Simons（吉姆·西蒙斯）", "size": 14, "bullet": True},
    {"text": "身份：世界级数学家、前NSA密码破译员", "size": 14, "bullet": True},
    {"text": "成立时间：1982年", "size": 14, "bullet": True},
    {"text": "管理资产：约 $1,300亿美元", "size": 14, "bullet": True},
    {"text": "", "size": 8},
    {"text": "大奖章基金 Medallion Fund", "size": 16, "color": COLOR_ACCENT, "bold": True, "space_after": 2},
    {"text": "被誉为「历史上最成功的对冲基金」", "size": 14, "color": COLOR_GOLD, "bullet": True},
    {"text": "平均年化回报率：费前 >60%，费后 ~40%", "size": 14, "color": COLOR_GREEN, "bullet": True},
    {"text": "自1988年来几乎每年正收益", "size": 14, "bullet": True},
    {"text": "已对外部投资者关闭，仅限员工投资", "size": 14, "bullet": True},
    {"text": "", "size": 8},
    {"text": "核心竞争力", "size": 16, "color": COLOR_ACCENT, "bold": True, "space_after": 2},
    {"text": "雇用数学家、物理学家（非传统金融人才）", "size": 14, "bullet": True},
    {"text": "使用复杂数学模型分析海量数据", "size": 14, "bullet": True},
    {"text": "极度保密，外界对其策略知之甚少", "size": 14, "bullet": True},
], title_size=16)

add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
         "为什么文艺复兴如此特别？", [
    {"text": "与传统对冲基金的关键区别", "size": 16, "color": COLOR_GOLD, "bold": True, "space_after": 2},
    {"text": "", "size": 4},
    {"text": "传统基金 → 雇用MBA、金融分析师", "size": 14, "bullet": True},
    {"text": "文艺复兴 → 雇用数学博士、物理学家、密码学家", "size": 14, "bullet": True, "color": COLOR_ACCENT},
    {"text": "", "size": 8},
    {"text": "传统基金 → 基于财报分析、行业研究", "size": 14, "bullet": True},
    {"text": "文艺复兴 → 纯数学模型、统计套利", "size": 14, "bullet": True, "color": COLOR_ACCENT},
    {"text": "", "size": 8},
    {"text": "传统基金 → 管理费2%、业绩报酬20%", "size": 14, "bullet": True},
    {"text": "文艺复兴 → 管理费5%、业绩报酬44%", "size": 14, "bullet": True, "color": COLOR_ACCENT},
    {"text": "（费率最高却仍供不应求）", "size": 13, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 8},
    {"text": "启示：在对冲基金行业，技术与人才是核心壁垒", "size": 14, "color": COLOR_GOLD},
    {"text": "Jim Simons 是陈-西蒙斯定理的共同创立者", "size": 13, "color": COLOR_LIGHT_GRAY},
], title_size=16)


# 索羅斯做空英鎊
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "经典交易  索罗斯做空英镑（1992年「黑色星期三」）", font_size=26, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(5), COLOR_ACCENT)

# 左：背景與邏輯
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5),
         "交易背景与逻辑", [
    {"text": "背景", "size": 16, "color": COLOR_ACCENT, "bold": True, "space_after": 2},
    {"text": "1990年英国加入欧洲汇率机制（ERM）", "size": 14, "bullet": True},
    {"text": "英镑被要求维持对德国马克的特定汇率区间", "size": 14, "bullet": True},
    {"text": "英国通胀率是德国的3倍，高利率伤害经济", "size": 14, "bullet": True},
    {"text": "", "size": 6},
    {"text": "索罗斯的判断", "size": 16, "color": COLOR_ACCENT, "bold": True, "space_after": 2},
    {"text": "英国加入ERM的汇率过高，基本面无法支撑", "size": 14, "bullet": True},
    {"text": "英国与德国的经济状况矛盾（一个需降息，一个需维持高息）", "size": 14, "bullet": True},
    {"text": "英国政府终将被迫退出ERM", "size": 14, "bullet": True},
    {"text": "", "size": 6},
    {"text": "交易策略", "size": 16, "color": COLOR_ACCENT, "bold": True, "space_after": 2},
    {"text": "大规模做空英镑，头寸从$15亿急增至$100亿", "size": 14, "color": COLOR_GOLD, "bullet": True},
], title_size=16)

# 右：結果
add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
         "交易结果", [
    {"text": "1992年9月16日 — 黑色星期三", "size": 16, "color": COLOR_RED, "bold": True, "space_after": 2},
    {"text": "", "size": 4},
    {"text": "英格兰银行拼命买入英镑 → 无效", "size": 14, "bullet": True},
    {"text": "紧急加息两次（10%→12%→15%）→ 仍无法阻止", "size": 14, "bullet": True},
    {"text": "英镑对马克暴跌15%，对美元暴跌25%", "size": 14, "bullet": True, "color": COLOR_RED},
    {"text": "英国被迫退出ERM", "size": 14, "bullet": True, "color": COLOR_RED},
    {"text": "", "size": 8},
    {"text": "收益：索罗斯获利超过 10亿英镑", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "量子基金从$150亿飙升至$220亿", "size": 14, "color": COLOR_GREEN},
    {"text": "英国财政部损失约33亿英镑", "size": 14, "color": COLOR_RED},
    {"text": "", "size": 8},
    {"text": "历史影响", "size": 16, "color": COLOR_ACCENT, "bold": True, "space_after": 2},
    {"text": "索罗斯被称为「击垮英格兰银行的人」", "size": 14, "bullet": True},
    {"text": "讽刺的是，脱离ERM后英国经济反而复苏", "size": 14, "bullet": True},
    {"text": "成为全球宏观策略的经典教材", "size": 14, "bullet": True, "color": COLOR_GOLD},
], title_size=16)


# LTCM 崩盤
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "惨痛教训  LTCM 崩盘（1998年）", font_size=28, color=COLOR_RED, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(5), COLOR_RED)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.8),
         "辉煌时期", [
    {"text": "1994年创立，团队包含两位诺贝尔经济学奖得主", "size": 14, "bullet": True},
    {"text": "核心策略：固定收益套利 / 相对价值策略", "size": 14, "bullet": True},
    {"text": "1994年回报 20% → 1995年 43% → 1996年 41%", "size": 14, "color": COLOR_GREEN, "bullet": True},
    {"text": "被华尔街视为「永远不会失败」的天才基金", "size": 14, "bullet": True},
], title_size=16)

add_card(slide, Inches(0.8), Inches(4.4), Inches(5.8), Inches(2.8),
         "崩盘过程", [
    {"text": "致命缺陷：杠杆比率高达 25:1", "size": 14, "color": COLOR_RED, "bullet": True},
    {"text": "1998年8月 — 俄罗斯债务违约引爆危机", "size": 14, "bullet": True},
    {"text": "全球「逃向安全资产」，流动性急剧枯竭", "size": 14, "bullet": True},
    {"text": "8月单月亏损44%（$21亿）", "size": 14, "color": COLOR_RED, "bullet": True},
    {"text": "全年亏损超过52%", "size": 14, "color": COLOR_RED, "bullet": True},
    {"text": "美联储协调14家机构注资$36.5亿救助", "size": 14, "bullet": True},
], title_size=16)

add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.9),
         "五大教训（投资者必知）", [
    {"text": "", "size": 4},
    {"text": "① 杠杆是双刃剑", "size": 16, "color": COLOR_RED, "bold": True},
    {"text": "可以放大收益，同样放大亏损", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "② 模型有局限性", "size": 16, "color": COLOR_RED, "bold": True},
    {"text": "再优秀的数学模型也无法预测极端事件（黑天鹅）", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "③ 流动性风险不可忽视", "size": 16, "color": COLOR_RED, "bold": True},
    {"text": "正常时期的相关性在危机时可能完全改变", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "④ 诺贝尔奖也救不了", "size": 16, "color": COLOR_RED, "bold": True},
    {"text": "学术荣誉不等于市场实战能力", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "⑤ 「大而不倒」风险", "size": 16, "color": COLOR_RED, "bold": True},
    {"text": "预示了2008年金融危机中系统性风险的概念", "size": 14, "color": COLOR_LIGHT_GRAY},
], title_size=16)


# 全球前十大
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "全球前十大对冲基金排名（2025-2026）", font_size=28,
             color=COLOR_DARK_BG, bold=True)

headers = ["排名", "基金名称", "管理资产（AUM）", "核心策略", "创始人/掌门人"]
rows = [
    ["1", "Millennium Management", "~$5,059亿", "多策略", "Israel Englander"],
    ["2", "Citadel", "~$3,970亿", "多策略", "Ken Griffin"],
    ["3", "Balyasny Asset Management", "~$2,480亿", "多策略", "Dmitry Balyasny"],
    ["4", "Bridgewater Associates", "~$1,500亿", "全球宏观", "Ray Dalio（创始）"],
    ["5", "D.E. Shaw", "~$850亿", "量化+基本面", "David E. Shaw"],
    ["6", "Two Sigma", "~$600亿", "量化/AI", "D. Siegel & J. Overdeck"],
    ["7", "AQR Capital", "大型", "量化", "Cliff Asness"],
    ["8", "Man Group", "大型", "多策略/量化", "Robyn Grew"],
    ["9", "Elliott Management", "大型", "事件驱动", "Paul Singer"],
    ["10", "TCI Fund Management", "大型", "集中持股", "Chris Hohn"],
]

add_table_slide(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.8),
                headers, rows, font_size=12,
                header_color=COLOR_DARK_BG)

add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
             "注：排名根据 HFR、HedgeFollow、InvestmentNews 等多个来源综合整理，不同统计口径可能有差异",
             font_size=11, color=RGBColor(0x66, 0x66, 0x66))


# ==================== 費用結構 ====================
make_section_title_slide("06  费用结构：2/20 模式",
                         "管理费 + 业绩报酬 + 高水位线机制")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "对冲基金费用结构详解", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

# 管理費
add_card(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(3.0),
         "管理费 2%", [
    {"text": "按管理资产总额每年收取", "size": 14, "bullet": True},
    {"text": "不论基金盈亏均需支付", "size": 14, "color": COLOR_RED, "bullet": True},
    {"text": "用途：员工薪酬、办公场地、行政管理", "size": 14, "bullet": True},
    {"text": "", "size": 6},
    {"text": "例：AUM $10亿", "size": 13, "color": COLOR_GOLD},
    {"text": "年管理费 = $2,000万", "size": 14, "color": COLOR_GREEN, "bold": True},
], title_size=16)

# 業績報酬
add_card(slide, Inches(4.8), Inches(1.3), Inches(3.7), Inches(3.0),
         "业绩报酬 20%", [
    {"text": "按基金利润的百分比收取", "size": 14, "bullet": True},
    {"text": "仅在基金获利时才收取", "size": 14, "color": COLOR_GREEN, "bullet": True},
    {"text": "部分设有门槛收益率（如8%）", "size": 14, "bullet": True},
    {"text": "", "size": 6},
    {"text": "例：获利 $5亿", "size": 13, "color": COLOR_GOLD},
    {"text": "业绩报酬 = $1亿", "size": 14, "color": COLOR_GREEN, "bold": True},
], title_size=16)

# 高水位線
add_card(slide, Inches(8.8), Inches(1.3), Inches(3.7), Inches(3.0),
         "高水位线机制", [
    {"text": "防止基金经理重复收费", "size": 14, "bullet": True},
    {"text": "净值须超过历史最高点", "size": 14, "bullet": True},
    {"text": "亏损后恢复的部分不收费", "size": 14, "color": COLOR_GOLD, "bullet": True},
    {"text": "", "size": 6},
    {"text": "保护投资者利益的关键机制", "size": 14, "color": COLOR_GREEN, "bold": True},
], title_size=16)

# 高水位線示例表
slide_fee2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide_fee2, RGBColor(0xF5, 0xF7, 0xFA))
add_text_box(slide_fee2, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "高水位线运作机制 — 实例说明", font_size=28,
             color=COLOR_DARK_BG, bold=True)

headers = ["时期", "基金净值", "高水位线", "是否收业绩报酬？", "说明"]
rows = [
    ["初始", "$100", "$100", "—", "投资起点"],
    ["第1年末", "$120", "$120", "是（对$20利润收取）", "盈利$20 × 20% = $4报酬"],
    ["第2年末", "$100", "$120", "否", "亏损，低于高水位线"],
    ["第3年末", "$110", "$120", "否", "虽然盈利但仍低于高水位线"],
    ["第4年末", "$140", "$140", "是（仅对$20收取）", "仅超出$120的$20收费"],
]
add_table_slide(slide_fee2, Inches(0.5), Inches(1.2), Inches(12.3), Inches(3.5),
                headers, rows, font_size=13,
                header_color=COLOR_DARK_BG)

# 費率趨勢
add_text_box(slide_fee2, Inches(0.8), Inches(5.0), Inches(10), Inches(0.5),
             "行业费率变化趋势", font_size=22, color=COLOR_DARK_BG, bold=True)

headers2 = ["时期", "平均管理费", "平均业绩报酬", "备注"]
rows2 = [
    ["传统模式", "2.00%", "20.00%", "经典「2/20」"],
    ["2019年行业平均", "1.50%", "19.00%", "投资者议价能力增强"],
    ["2025年行业平均", "<1.50%", "~15%", "费率持续压缩"],
    ["顶级基金（如Medallion）", "5%", "44%", "需求远超供给，反而提价"],
]
add_table_slide(slide_fee2, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.8),
                headers2, rows2, font_size=12,
                header_color=COLOR_DARK_BG)


# ==================== 風險與監管 ====================
make_section_title_slide("07  风险全景与全球监管",
                         "六大风险类型 · 美欧亚监管框架 · 投资者保护")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "对冲基金六大风险类型", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

risks = [
    ("市场风险", "Market Risk", "资产价格波动造成的损失\n极端市场中相关性可能突变", COLOR_RED),
    ("杠杆风险", "Leverage Risk", "对冲基金的杠杆没有法定上限\nLTCM的25:1杠杆是致命因素", COLOR_RED),
    ("流动性风险", "Liquidity Risk", "投资者集中赎回时\n基金可能被迫低价抛售", COLOR_ORANGE),
    ("交易对手风险", "Counterparty", "交易对手（银行等）可能违约\n场外衍生品交易中尤为突出", COLOR_ORANGE),
    ("操作风险", "Operational", "估值不当、内控薄弱\n麦道夫骗局即为极端案例", COLOR_GOLD),
    ("模型风险", "Model Risk", "量化模型在未知市场可能失效\n「黑天鹅」难以被模型捕捉", COLOR_GOLD),
]

for i, (name, en, desc, color) in enumerate(risks):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + Inches(4.1) * col
    top = Inches(1.3) + Inches(3.0) * row

    add_shape_bg(slide, left, top, Inches(3.8), Inches(2.7), COLOR_MID_BG)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.5), Inches(0.35),
                 f"{name}  {en}", font_size=16, color=color, bold=True)
    add_divider_line(slide, left + Inches(0.15), top + Inches(0.5), Inches(3.5), color, Pt(1))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.65), Inches(3.5), Inches(1.8),
                 desc, font_size=14, color=COLOR_LIGHT_GRAY)


# 全球監管
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "全球监管框架与投资者保护", font_size=28, color=COLOR_DARK_BG, bold=True)

headers = ["地区", "监管机构", "核心要求", "关键法规"]
rows = [
    ["美国", "SEC / CFTC", "AUM>$1.5亿须注册\n定期提交Form PF报告", "Dodd-Frank法案（2010）"],
    ["欧洲", "各国监管+ESMA", "基金经理须在欧盟注册\n报告杠杆、风险、流动性", "AIFMD指令（2011）"],
    ["香港", "SFC证监会", "持牌制度\n第9类（资产管理）牌照", "《证券及期货条例》"],
    ["新加坡", "MAS金管局", "基金公司须持CMS牌照", "《证券期货法》"],
    ["开曼群岛", "CIMA", "全球最受欢迎注册地\n监管逐步加强", "SIBL / MFL"],
]
add_table_slide(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.5),
                headers, rows, font_size=12, header_color=COLOR_DARK_BG)

# 投資者保護
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.5),
             "投资者保护机制", font_size=22, color=COLOR_DARK_BG, bold=True)

headers_p = ["保护机制", "说明"]
rows_p = [
    ["尽职调查 Due Diligence", "审查基金文件、管理人履历、业绩历史、风控体系"],
    ["独立托管人 Custodian", "资产由独立机构托管，防止挪用"],
    ["独立审计", "年度由四大会计师事务所审计"],
    ["独立行政管理人", "独立计算净值，防止人为操控"],
    ["投资者信函 Side Letters", "大型投资者可协商额外保护条款"],
]
add_table_slide(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.8),
                headers_p, rows_p, font_size=12, header_color=COLOR_DARK_BG)


# ==================== 市場規模與趨勢 ====================
make_section_title_slide("08  全球市场规模与趋势",
                         "$5.15万亿 — 历史性里程碑")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "全球对冲基金市场核心数据", font_size=28, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

# 核心數字
data_cards = [
    ("$5.15万亿", "2025年末管理资产", "HFR数据，历史新高", COLOR_GREEN),
    ("$8.83万亿", "2031年预测规模", "CAGR 9.12%", COLOR_ACCENT),
    ("11.2%", "2025年平均回报率", "连续两年双位数", COLOR_GOLD),
    ("$6,428亿", "2025年资本增长", "创历史纪录", COLOR_GREEN),
]

for i, (number, label, note, color) in enumerate(data_cards):
    left = Inches(0.6) + Inches(3.15) * i
    add_shape_bg(slide, left, Inches(1.3), Inches(2.9), Inches(2.0), COLOR_MID_BG)
    add_text_box(slide, left + Inches(0.1), Inches(1.4), Inches(2.7), Inches(0.8),
                 number, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.2), Inches(2.7), Inches(0.4),
                 label, font_size=14, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.6), Inches(2.7), Inches(0.4),
                 note, font_size=12, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 六大趨勢
add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(0.5),
             "六大行业趋势", font_size=22, color=COLOR_GOLD, bold=True)

trends = [
    ("多策略基金崛起", "2025年首次超越股票策略\n成为最大策略类别"),
    ("量化与AI革命", "量化策略CAGR 11.63%\nAI深度融入投资决策"),
    ("费率持续压缩", "平均管理费降至1.5%以下\n顶级基金反而提价"),
    ("亚太市场加速", "配置占比24%→30%\n中国、印度、日本为重点"),
    ("ESG永续投资", "越来越多基金纳入ESG\n监管要求增加相关披露"),
    ("行业集中化", "551家大基金管理\n行业86%的资产"),
]

for i, (title, desc) in enumerate(trends):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + Inches(4.1) * col
    top = Inches(4.2) + Inches(1.55) * row

    add_shape_bg(slide, left, top, Inches(3.8), Inches(1.35), COLOR_MID_BG)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), Inches(3.5), Inches(0.35),
                 title, font_size=15, color=COLOR_ACCENT, bold=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.45), Inches(3.5), Inches(0.8),
                 desc, font_size=12, color=COLOR_LIGHT_GRAY)


# ==================== 中國投資者參與渠道 ====================
make_section_title_slide("09  中国投资者参与渠道",
                         "QDII · QDLP · QDIE · 直接投资 — 四大通道解析")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "中国投资者参与海外对冲基金的四大渠道", font_size=26, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(5), COLOR_ACCENT)

channels = [
    ("QDII", "合格境内机构投资者", [
        "经证监会批准的基金公司、保险、券商",
        "门槛较低，可通过银行购买",
        "受额度限制，多数不直接投资对冲基金",
    ]),
    ("QDLP", "合格境内有限合伙人", [
        "试点城市：上海、天津等",
        "允许境外基金管理人在中国设立LP企业",
        "可直接投资海外对冲基金",
        "门槛较高，面向机构和高净值个人",
    ]),
    ("QDIE", "合格境内投资企业", [
        "试点城市：深圳前海",
        "与QDLP类似的投资通道",
        "可投对冲基金、PE、VC、房地产基金等",
    ]),
    ("直接海外投资", "通过海外银行/家族办公室", [
        "通过海外银行账户直接投资",
        "需符合外汇管理规定",
        "每人每年$5万便利额度",
        "大额资金需合规通道",
    ]),
]

for i, (name, full_name, points) in enumerate(channels):
    left = Inches(0.5) + Inches(3.15) * i
    lines = [{"text": full_name, "size": 13, "color": COLOR_ACCENT}]
    lines += [{"text": p, "size": 12, "bullet": True} for p in points]
    add_card(slide, left, Inches(1.3), Inches(2.9), Inches(3.5),
             name, lines, title_size=18, body_size=12)

# 門檻表
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.5),
             "投资门槛对照", font_size=20, color=COLOR_GOLD, bold=True)

add_shape_bg(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), COLOR_MID_BG)
threshold_lines = [
    {"text": "QDII基金产品：1,000-10,000人民币 → 适合一般投资者", "size": 14, "bullet": True},
    {"text": "QDLP/QDIE：100万-300万人民币 → 适合合格投资者", "size": 14, "bullet": True},
    {"text": "直接投资海外对冲基金：$100万-$500万美元 → 适合超高净值/机构", "size": 14, "bullet": True},
    {"text": "通过FOF（母基金）：$25万-$100万美元 → 适合高净值投资者", "size": 14, "bullet": True},
]
add_multi_text(slide, Inches(0.8), Inches(5.55), Inches(11.5), Inches(1.5),
               threshold_lines, default_size=14)


# 中國投資者風險
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "中国投资者需特别注意的风险", font_size=28, color=COLOR_RED, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(5), COLOR_RED)

cn_risks = [
    ("外汇风险", "投资美元计价基金\n人民币汇率波动影响实际收益", COLOR_RED),
    ("监管风险", "中国外汇管制政策可能变化\n影响资金跨境流动", COLOR_RED),
    ("信息不对称", "海外基金信息披露制度不同\n语言和文化障碍", COLOR_ORANGE),
    ("法律管辖风险", "多数基金注册在开曼等离岸地区\n法律纠纷处理困难", COLOR_ORANGE),
    ("流动性风险", "锁定期和赎回限制\n急需资金时可能无法退出", COLOR_GOLD),
    ("税务合规风险", "CRS共同申报准则\n海外投资税务申报要求", COLOR_GOLD),
]

for i, (name, desc, color) in enumerate(cn_risks):
    col = i % 3
    row = i // 3
    left = Inches(0.6) + Inches(4.1) * col
    top = Inches(1.3) + Inches(3.0) * row

    add_shape_bg(slide, left, top, Inches(3.8), Inches(2.7), COLOR_MID_BG)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(3.5), Inches(0.4),
                 name, font_size=18, color=color, bold=True)
    add_divider_line(slide, left + Inches(0.15), top + Inches(0.5), Inches(3.5), color, Pt(1))
    add_text_box(slide, left + Inches(0.15), top + Inches(0.65), Inches(3.5), Inches(1.8),
                 desc, font_size=15, color=COLOR_LIGHT_GRAY)


# ==================== 術語表 ====================
make_section_title_slide("10  核心术语表与总结",
                         "20个必知专业术语 · 投资建议")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RGBColor(0xF5, 0xF7, 0xFA))
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "核心术语表（中英对照）", font_size=28, color=COLOR_DARK_BG, bold=True)

headers = ["英文", "中文", "简要说明"]
rows = [
    ["Hedge Fund", "对冲基金", "使用多种策略追求绝对收益的私募基金"],
    ["Alpha", "超额收益", "基金经理选股能力带来的超额回报"],
    ["Beta", "市场风险", "基金随市场整体波动的程度"],
    ["AUM", "管理资产规模", "Assets Under Management"],
    ["Long / Short", "做多 / 做空", "买入看好资产 / 卖出借入资产"],
    ["Leverage", "杠杆", "借钱投资以放大收益（和风险）"],
    ["High-Water Mark", "高水位线", "业绩报酬计算基准"],
    ["Lock-up Period", "锁定期", "投资者不得赎回的期间"],
    ["Drawdown", "回撤", "从最高点下跌的幅度"],
    ["Fund of Funds", "母基金（FOF）", "投资于多支对冲基金的基金"],
    ["Prime Broker", "主经纪商", "为对冲基金提供融资借券等服务"],
    ["Sharpe Ratio", "夏普比率", "风险调整后收益指标"],
]
add_table_slide(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0),
                headers, rows, font_size=12, header_color=COLOR_DARK_BG)


# ==================== 總結頁 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
             "总结与投资建议", font_size=32, color=COLOR_GOLD, bold=True)
add_divider_line(slide, Inches(0.8), Inches(1.0), Inches(4), COLOR_ACCENT)

add_card(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.8),
         "对冲基金的核心价值", [
    {"text": "", "size": 4},
    {"text": "① 绝对收益能力", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "在牛市和熊市中均可追求正回报", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "② 分散投资风险", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "与传统股票债券低相关，优化组合效率", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "③ 策略多元化", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "六大策略适应不同市场环境", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "④ 全球配置视角", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "跨市场、跨资产、跨地区投资", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "⑤ 行业持续增长", "size": 16, "color": COLOR_ACCENT, "bold": True},
    {"text": "管理资产从$390亿增长至$5.15万亿", "size": 14, "color": COLOR_GREEN},
], title_size=16)

add_card(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.8),
         "给中国投资者的建议", [
    {"text": "", "size": 4},
    {"text": "✓ 充分了解后再投资", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "不盲目追捧，先理解策略逻辑和风险", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "✓ 合规渠道参与", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "通过QDII、QDLP、FOF等合规通道", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "✓ 分散配置", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "对冲基金作为资产组合的一部分，占比10-20%", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "✓ 重视尽职调查", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "审查基金经理背景、历史业绩、风控体系", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "✓ 注意流动性安排", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "根据自身资金需求选择合适锁定期", "size": 14, "color": COLOR_LIGHT_GRAY},
    {"text": "", "size": 6},
    {"text": "✓ 关注外汇与税务风险", "size": 16, "color": COLOR_GREEN, "bold": True},
    {"text": "人民币汇率波动、CRS合规申报", "size": 14, "color": COLOR_LIGHT_GRAY},
], title_size=16)


# ==================== 封底 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, COLOR_DARK_BG)

add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), COLOR_GOLD)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1),
             "感谢阅读", font_size=48, color=COLOR_WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
             "THANK YOU", font_size=28, color=COLOR_GOLD,
             alignment=PP_ALIGN.CENTER)

add_divider_line(slide, Inches(4.5), Inches(4.1), Inches(4.3), COLOR_ACCENT, Pt(2))

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
             "海外对冲基金  |  投资者教育系列", font_size=18,
             color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8),
             "数据来源：HFR · Preqin · Bloomberg · J.P. Morgan\n"
             "Institutional Investor · Man Group · Wellington · Barclays",
             font_size=14, color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             "免责声明：本材料仅供教育目的，不构成任何投资建议。投资有风险，入市须谨慎。",
             font_size=12, color=RGBColor(0x80, 0x88, 0x98), alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, Pt(4), COLOR_GOLD)


# ==================== 保存 ====================
output_path = "/Users/jamie/Desktop/Claude-股票量化數據研究網/对冲基金/海外对冲基金_投资者教育.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
print(f"共 {len(prs.slides)} 页")
